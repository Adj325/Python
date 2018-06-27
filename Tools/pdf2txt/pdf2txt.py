import os
try:
    import pptx
except:
    os.system('pip3 install python-pptx')  
    import pptx
        
ppt_types = ['ppt', 'pptx', 'pptm']

def pdf_to_txt(pdf_path):

    txt_content =  ''
    
    # 打开pdf
    try:
        ppt_file = pptx.Presentation(pdf_path)
    except:
        print('警告：\r\t无法处理"{}"'.format(pdf_path))
        print('\t此文件可能不是ppt文件\n')
        return

    # 获取页数
    page_num = len(ppt_file.slides)
    for n in range(page_num):
                
        for shap in ppt_file.slides[n].shapes:
            try:
                txt_content += shap.text+'\n'
                #print(shap.text)
            except:
                pass
    # 美化内容
    txt_content = beautfy(txt_content)
    pdf_name = pdf_path.split('\\')[-1]
    try:
        with open(pdf_name+'.txt', 'w', encoding='gb18030') as f:
            f.write(txt_content)
    except:
        with open(pdf_name+'.txt', 'w', encoding='utf-8') as f:
            f.write(txt_content)
    print('\n提示：\n\t"{}" 处理成功！\n\t"{}.txt" 保存在当前目录下\n'.format(pdf_path, pdf_name))

def beautfy(content):

    flag = False

    while True:
        if flag:
            return content
        flag = True
        
        # 统一换行符号为\r\n（兼容windows）
        content = content.replace('\r\n', '\n')
        content = content.replace('\n', '\r\n')
        while '\r\n\r\n\r\n' in content:
            content = content.replace('\r\n\r\n\r\n', '\r\n\r\n')
            flag = False

        # 去除多余空格
        while '  ' in content:
            content = content.replace('  ', '')
            flag = False

        # 空格行
        while '\r\n  \r\n' in content:
            content = content.replace('\r\n  \r\n', '\r\n\r\n')
            flag = False

if __name__ == '__main__':
    print('程序功能: 提取ppt的文本内容')
    print('PPT 类型: '+' '.join(ppt_types))
    print()
    choice = input('回车 - 处理当前目录所有ppt\n其它 - 处理特定ppt(文件若在当前目录，只需要填写文件名，否则需要填写完整路径)\n\n选择: ')
    if choice == '':
        print('处理当前目录所有ppt\n')
        pdf_names = [f for f in os.listdir() if f.split('.')[-1].lower() in ppt_types]

        for pdf_name in pdf_names:
            pdf_to_txt(pdf_name)
    else:
        print('Ctrl+C退出\n')
        while True:
            pdf_path = input('pdf路径:')
            pdf_to_txt(pdf_path)
